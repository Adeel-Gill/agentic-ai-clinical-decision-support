#!/usr/bin/env python3
"""
Build the v1.0 defense deck using the Superior University Design System v1.0
(see Presentation_Design_Guide.md).

Every content slide carries a real, natively-drawn diagram (no placeholders):
rounded nodes, consistent arrows, brand colours, flat fills — per the guide's
diagram rules.

Output:  Defense_Presentation_v1.pptx  -- white, red-accent, academic 16:9 deck.
Run:     python build_presentation_v1.py
Requires: python-pptx, Pillow  (logo at ../07_Thesis/Images/superior_logo.png)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
LOGO = os.path.join(HERE, "..", "07_Thesis", "Images", "superior_logo.png")

# ---- Design System v1.0 palette -----------------------------------------
RED      = RGBColor(0xD7, 0x29, 0x24)
ROSE     = RGBColor(0xCB, 0x72, 0x5E)
PURPLE   = RGBColor(0x6C, 0x1C, 0x74)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
INK      = RGBColor(0x22, 0x22, 0x22)
SLATE    = RGBColor(0x55, 0x55, 0x55)
MUTE     = RGBColor(0x8A, 0x8A, 0x8A)
SURFACE  = RGBColor(0xF5, 0xF5, 0xF5)
HAIRLINE = RGBColor(0xEA, 0xEA, 0xEA)
GREEN    = RGBColor(0x3F, 0x8A, 0x5B)
ORANGE   = RGBColor(0xC8, 0x78, 0x1E)
PTINT    = RGBColor(0xF3, 0xE9, 0xF5)   # pale purple highlight
GTINT    = RGBColor(0xE9, 0xF2, 0xEC)
OTINT    = RGBColor(0xF7, 0xEE, 0xE1)
DARK     = RGBColor(0x1A, 0x11, 0x13)
DARKSUB  = RGBColor(0xCC, 0xC4, 0xC6)

FONT = "Aptos"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
SHORT_TITLE = "Agentic AI  ·  Clinical Decision Support"
UNIV        = "The Superior University"

# ---- figure card geometry (inches) --------------------------------------
CARD_X, CARD_Y, CARD_W, CARD_H = 8.35, 2.05, 4.28, 4.35
DX = CARD_X + 0.30      # diagram region left
DY = CARD_Y + 1.15      # diagram region top
DW = CARD_W - 0.60      # diagram region width  (3.68)
DH = CARD_H - 1.45      # diagram region height (2.90)
CXc = DX + DW / 2       # horizontal centre (10.49)


# ==========================================================================
# Primitives
# ==========================================================================
def _no_shadow(shape):
    shape.shadow.inherit = False


def _soft_shadow(shape):
    sp = shape._element.spPr
    eff = sp.find(qn('a:effectLst'))
    if eff is None:
        eff = sp.makeelement(qn('a:effectLst'), {}); sp.append(eff)
    sh = eff.makeelement(qn('a:outerShdw'),
                         {'blurRad': '90000', 'dist': '38000',
                          'dir': '5400000', 'rotWithShape': '0'})
    c = sh.makeelement(qn('a:srgbClr'), {'val': '000000'})
    a = c.makeelement(qn('a:alpha'), {'val': '11000'}); c.append(a); sh.append(c)
    eff.append(sh)


def add_rect(slide, x, y, w, h, color, line=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = line_w or Pt(1)
    _no_shadow(shp)
    return shp


def add_round(slide, x, y, w, h, fill, line=None, line_w=None, radius=0.10, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = line_w or Pt(1)
    _no_shadow(shp)
    if shadow:
        _soft_shadow(shp)
    return shp


def set_text(shape, paras, align=PP_ALIGN.CENTER, padl=3, anchor=MSO_ANCHOR.MIDDLE, ls=1.0):
    """paras = list of paragraphs; each paragraph = list of (text,size,color,bold)."""
    tf = shape.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(padl); tf.margin_right = Pt(3)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    tf.vertical_anchor = anchor
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = ls
        for (t, sz, col, b) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = b; r.font.name = FONT


def textbox(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space_after=0, ls=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space_after:
            p.space_after = Pt(space_after)
        if ls:
            p.line_spacing = ls
        for (t, sz, col, b, it) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = b; r.font.italic = it; r.font.name = FONT
    return tb


def node(slide, x, y, w, h, text, fill=WHITE, edge=HAIRLINE, tcolor=INK,
         size=10, bold=False, ew=1.0, rad=0.14, align=PP_ALIGN.CENTER, shadow=False):
    shp = add_round(slide, Inches(x), Inches(y), Inches(w), Inches(h), fill,
                    line=edge, line_w=Pt(ew), radius=rad, shadow=shadow)
    lines = text if isinstance(text, list) else [text]
    set_text(shp, [[(t, size, tcolor, bold)] for t in lines], align=align, padl=5)
    return shp


def _line_xml(cxn, color, w, arrow, dash):
    cxn.line.color.rgb = color; cxn.line.width = Pt(w); _no_shadow(cxn)
    ln = cxn.line._get_or_add_ln()
    if dash:
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    if arrow:
        ln.append(ln.makeelement(qn('a:tailEnd'),
                                 {'type': 'triangle', 'w': 'med', 'len': 'med'}))


def arrow(slide, x1, y1, x2, y2, color=SLATE, w=1.4, dash=False, head=True):
    cxn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                     Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    _line_xml(cxn, color, w, head, dash)
    return cxn


def oval(slide, x, y, d, edge, ew=2.0, fill=None):
    o = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    if fill is None:
        o.fill.background()
    else:
        o.fill.solid(); o.fill.fore_color.rgb = fill
    o.line.color.rgb = edge; o.line.width = Pt(ew); _no_shadow(o)
    return o


def badge(slide, x, y, d, color=GREEN, mark='✓'):
    o = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    o.fill.solid(); o.fill.fore_color.rgb = color; o.line.fill.background(); _no_shadow(o)
    set_text(o, [[(mark, 11, WHITE, True)]])
    return o


def diamond(slide, x, y, w, h, text, edge=PURPLE, tcolor=INK, size=9.5):
    d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(x), Inches(y), Inches(w), Inches(h))
    d.fill.solid(); d.fill.fore_color.rgb = PTINT
    d.line.color.rgb = edge; d.line.width = Pt(1.25); _no_shadow(d)
    set_text(d, [[(text, size, tcolor, False)]])
    return d


def caption_mark(slide, x, y, w, sym, color):
    """A ✓ / ◐ / ✗ symbol centred in a matrix cell."""
    textbox(slide, Inches(x), Inches(y), Inches(w), Inches(0.34),
            [[(sym, 12, color, True, False)]], align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)


# ==========================================================================
# Diagram functions  (draw inside DX,DY,DW,DH)
# ==========================================================================
def fig_bottleneck(s):
    nw, nh = 1.72, 0.44
    x1, x2 = DX, DX + DW - nw
    y1, y2 = DY, DY + 0.54
    for (x, y, t) in [(x1, y1, "Vitals"), (x2, y1, "Labs"),
                      (x1, y2, "Notes"), (x2, y2, "Imaging")]:
        node(s, x, y, nw, nh, t, fill=SURFACE, size=10)
    bx, by, bw, bh = CXc - 1.0, DY + 1.35, 2.0, 0.44
    node(s, bx, by, bw, bh, "Manual synthesis", fill=WHITE, size=10)
    cx, cy, cw, ch = CXc - 0.95, DY + 2.28, 1.9, 0.5
    node(s, cx, cy, cw, ch, "Clinician", edge=RED, ew=1.5, bold=True, size=11)
    for (x, y) in [(x1 + nw / 2, y1 + nh), (x2 + nw / 2, y1 + nh),
                   (x1 + nw / 2, y2 + nh), (x2 + nw / 2, y2 + nh)]:
        arrow(s, x, y, CXc, by)
    arrow(s, CXc, by + bh, CXc, cy)


def fig_passive_agentic(s):
    textbox(s, Inches(DX), Inches(DY - 0.02), Inches(DW), Inches(0.24),
            [[("PASSIVE  LLM", 9, MUTE, True, False)]])
    y = DY + 0.24
    w, h = 1.02, 0.42
    xs = [DX, DX + 1.33, DX + 2.66]
    for x, t in zip(xs, ["Prompt", "Black box", "Response"]):
        node(s, x, y, w, h, t, fill=SURFACE, size=9.5)
    arrow(s, xs[0] + w, y + h / 2, xs[1], y + h / 2)
    arrow(s, xs[1] + w, y + h / 2, xs[2], y + h / 2)

    textbox(s, Inches(DX), Inches(DY + 1.00), Inches(DW), Inches(0.24),
            [[("AGENTIC  LOOP", 9, RED, True, False)]])
    y2 = DY + 1.26
    w2 = 0.85
    gap = (DW - 4 * w2) / 3
    xs2 = [DX + i * (w2 + gap) for i in range(4)]
    for i, (x, t) in enumerate(zip(xs2, ["Perceive", "Memory", "Reason", "Act"])):
        node(s, x, y2, w2, 0.44, t, fill=SURFACE, size=9.5,
             edge=(RED if t == "Reason" else HAIRLINE))
    for i in range(3):
        arrow(s, xs2[i] + w2, y2 + 0.22, xs2[i + 1], y2 + 0.22)
    # feedback loop
    ay = y2 + 0.72
    arrow(s, xs2[3] + w2 / 2, y2 + 0.44, xs2[3] + w2 / 2, ay, head=False, dash=True, color=MUTE)
    arrow(s, xs2[3] + w2 / 2, ay, xs2[0] + w2 / 2, ay, dash=True, color=MUTE)
    arrow(s, xs2[0] + w2 / 2, ay, xs2[0] + w2 / 2, y2 + 0.44, head=True, dash=True, color=MUTE)
    textbox(s, Inches(DX), Inches(ay + 0.06), Inches(DW), Inches(0.22),
            [[("environmental feedback", 8.5, MUTE, False, True)]], align=PP_ALIGN.CENTER)


def fig_pillars(s):
    labels = ["Lack of context", "Opaque reasoning", "Fragmented systems"]
    w, h = 1.06, 1.75
    gap = (DW - 3 * w) / 2
    y = DY + 0.30
    for i, t in enumerate(labels):
        x = DX + i * (w + gap)
        node(s, x, y, w, h, "", fill=SURFACE, edge=HAIRLINE, rad=0.10)
        add_rect(s, Inches(x), Inches(y), Inches(w), Pt(4), ORANGE)
        badge(s, x + w / 2 - 0.16, y - 0.16, 0.32, color=ORANGE, mark='✕')
        set_text(add_round(s, Inches(x), Inches(y), Inches(w), Inches(h), None, radius=0.10),
                 [[(t, 10.5, INK, False)]], align=PP_ALIGN.CENTER)
    add_rect(s, Inches(DX), Inches(y + h + 0.06), Inches(DW), Pt(1.5), HAIRLINE)


def _matrix(s, rows, cols, marks):
    c0 = 1.28
    cw = (DW - c0) / len(cols)
    rh = 0.40
    x0, y0 = DX, DY
    total_h = rh * (len(rows) + 1)
    # outer rounded border
    add_round(s, Inches(x0), Inches(y0), Inches(DW), Inches(total_h), None,
              line=HAIRLINE, line_w=Pt(1), radius=0.05)
    # header
    add_rect(s, Inches(x0), Inches(y0), Inches(DW), Inches(rh), RED)
    textbox(s, Inches(x0 + 0.08), Inches(y0), Inches(c0 - 0.1), Inches(rh),
            [[("Framework", 9.5, WHITE, True, False)]], anchor=MSO_ANCHOR.MIDDLE)
    for j, c in enumerate(cols):
        textbox(s, Inches(x0 + c0 + j * cw), Inches(y0), Inches(cw), Inches(rh),
                [[(c, 9, WHITE, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # rows
    for i, rname in enumerate(rows):
        ry = y0 + rh * (i + 1)
        if i % 2 == 1:
            add_rect(s, Inches(x0), Inches(ry), Inches(DW), Inches(rh), SURFACE)
        textbox(s, Inches(x0 + 0.1), Inches(ry), Inches(c0 - 0.1), Inches(rh),
                [[(rname, 9.5, INK, False, False)]], anchor=MSO_ANCHOR.MIDDLE)
        for j in range(len(cols)):
            sym, col = marks[i][j]
            caption_mark(s, x0 + c0 + j * cw, ry + 0.02, cw, sym, col)
    for i in range(len(rows)):
        ry = y0 + rh * (i + 1)
        add_rect(s, Inches(x0), Inches(ry), Inches(DW), Pt(0.75), HAIRLINE)


def fig_frameworks(s):
    Y, R, X = ('✓', GREEN), ('◐', ROSE), ('✕', MUTE)
    rows = ["ReAct", "AutoGen", "CAMEL", "MetaGPT"]
    cols = ["Memory", "SOPs", "Oversight"]
    marks = [[X, X, X], [R, X, R], [R, X, X], [R, Y, R]]
    _matrix(s, rows, cols, marks)
    textbox(s, Inches(DX), Inches(DY + 0.40 * 5 + 0.08), Inches(DW), Inches(0.3),
            [[("✓ met   ◐ partial   ✕ absent", 8.5, MUTE, False, True)]],
            align=PP_ALIGN.CENTER)


def fig_rag(s):
    tw = 1.72
    node(s, DX, DY, tw, 0.46, ["Medical", "knowledge"], fill=SURFACE, size=9.5)
    node(s, DX + DW - tw, DY, tw, 0.46, ["EHR", "databases"], fill=SURFACE, size=9.5)
    ry = DY + 0.80
    node(s, CXc - 1.0, ry, 2.0, 0.46, "Retriever", size=10)
    gy = ry + 0.78
    node(s, CXc - 1.0, gy, 2.0, 0.5, "Generator LLM", edge=RED, ew=1.5, bold=True, size=10.5)
    ay = gy + 0.82
    node(s, CXc - 1.2, ay, 2.4, 0.46, "Grounded recommendation", size=9.5)
    arrow(s, DX + tw / 2, DY + 0.46, CXc - 0.5, ry)
    arrow(s, DX + DW - tw / 2, DY + 0.46, CXc + 0.5, ry)
    arrow(s, CXc, ry + 0.46, CXc, gy)
    arrow(s, CXc, gy + 0.5, CXc, ay)


def fig_venn(s):
    d = 1.75
    ax, ay = DX + 0.10, DY
    bx, by = DX + DW - d - 0.10, DY
    cx, cy = DX + (DW - d) / 2, DY + 0.92
    oval(s, ax, ay, d, RED)
    oval(s, bx, by, d, ROSE)
    oval(s, cx, cy, d, PURPLE)
    textbox(s, Inches(ax - 0.05), Inches(ay + 0.22), Inches(1.15), Inches(0.6),
            [[("Longitudinal ICU eval", 8.5, RED, True, False)]], align=PP_ALIGN.CENTER)
    textbox(s, Inches(bx + d - 1.10), Inches(by + 0.22), Inches(1.15), Inches(0.6),
            [[("Patient-timeline RAG", 8.5, ROSE, True, False)]], align=PP_ALIGN.CENTER)
    textbox(s, Inches(cx + d / 2 - 0.6), Inches(cy + d - 0.55), Inches(1.2), Inches(0.5),
            [[("Faithful auditability", 8.5, PURPLE, True, False)]], align=PP_ALIGN.CENTER)
    textbox(s, Inches(CXc - 0.62), Inches(DY + 1.02), Inches(1.24), Inches(0.6),
            [[("Proposed", 9, PURPLE, True, False)], [("framework", 9, PURPLE, True, False)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def fig_objectives(s):
    ccx, ccy, cw, ch = CXc - 0.98, DY + 1.02, 1.96, 0.74
    top = (CXc - 0.85, DY, 1.7, 0.5, "Memory · RAG · ReAct")
    bl = (DX, DY + 2.18, 1.62, 0.62, ["Specialized", "clinical agents"])
    br = (DX + DW - 1.62, DY + 2.18, 1.62, 0.62, ["Trustworthy AI", "+ HITL"])
    cpt = (ccx + cw / 2, ccy + ch / 2)
    for (x, y, w, h, _) in [top, bl, br]:
        arrow(s, cpt[0], cpt[1], x + w / 2, y + h / 2, head=False, color=HAIRLINE, w=1.5)
    node(s, ccx, ccy, cw, ch, ["Agentic CDSS", "framework"], edge=RED, ew=1.6,
         bold=True, size=11, fill=SURFACE)
    for (x, y, w, h, t) in [top, bl, br]:
        node(s, x, y, w, h, t, fill=WHITE, size=9.5)


def fig_rqs(s):
    items = [("RQ1", "Architectural limits"), ("RQ2", "Reasoning + physician control"),
             ("RQ3", "Agent collaboration & XAI"), ("RQ4", "RAG reliability & failure"),
             ("RQ5", "MIMIC-IV instantiation")]
    w = 3.34
    x = DX + (DW - w) / 2
    h = 0.38
    step = 0.545
    for i, (tag, lbl) in enumerate(items):
        y = DY + i * step
        shp = add_round(s, Inches(x), Inches(y), Inches(w), Inches(h), SURFACE,
                        line=HAIRLINE, line_w=Pt(1), radius=0.20)
        set_text(shp, [[(tag + "   ", 9.5, RED, True), (lbl, 9.5, INK, False)]],
                 align=PP_ALIGN.LEFT, padl=9)
        if i < len(items) - 1:
            arrow(s, x + w / 2, y + h, x + w / 2, y + step, color=MUTE, w=1.2)


def fig_methodology(s):
    ph = [("PHASE 1", "Literature review"), ("PHASE 2", "Taxonomy & gap analysis"),
          ("PHASE 3", "Framework design"), ("PHASE 4", "Prototype & evaluation plan")]
    w = 3.34
    x = DX + (DW - w) / 2
    h = 0.52
    step = 0.735
    for i, (tag, lbl) in enumerate(ph):
        y = DY + i * step
        node(s, x, y, w, h, "", fill=SURFACE, rad=0.10)
        add_rect(s, Inches(x), Inches(y), Pt(4), Inches(h), RED)
        textbox(s, Inches(x + 0.16), Inches(y + 0.06), Inches(w - 0.3), Inches(0.4),
                [[(tag, 8, RED, True, False)], [(lbl, 10.5, INK, False, False)]], ls=1.02)
        if i < 3:
            arrow(s, x + w / 2, y + h, x + w / 2, y + step, color=MUTE, w=1.2)


def fig_architecture(s):
    bands = ["Perception", "Memory", "Reasoning", "Action"]
    bw = 2.80
    bh = 0.50
    step = 0.635
    x = DX
    for i, t in enumerate(bands):
        y = DY + i * step
        node(s, x, y, bw, bh, t, fill=SURFACE, size=10.5,
             edge=(RED if t == "Reasoning" else HAIRLINE))
        if i < 3:
            arrow(s, x + bw / 2, y + bh, x + bw / 2, y + step, color=MUTE, w=1.2)
    # trust pillar
    px = DX + bw + 0.14
    pw = DW - bw - 0.14
    ph = step * 3 + bh
    p = add_round(s, Inches(px), Inches(DY), Inches(pw), Inches(ph), PTINT,
                  line=PURPLE, line_w=Pt(1), radius=0.06)
    set_text(p, [[("Trust", 10, PURPLE, True)], [("&", 10, PURPLE, True)],
                 [("verification", 10, PURPLE, True)]], ls=1.15)


def fig_hub(s):
    hcx, hcy, hw, hh = CXc - 0.85, DY + 1.05, 1.7, 0.6
    sats = [(DX, DY, "Monitoring"), (DX + DW - 1.5, DY, "Diagnosis"),
            (DX, DY + 2.30, "Treatment"), (DX + DW - 1.5, DY + 2.30, ["Explanation", "/ verify"])]
    hc = (hcx + hw / 2, hcy + hh / 2)
    for (x, y, _) in sats:
        arrow(s, hc[0], hc[1], x + 0.75, y + 0.23, head=False, color=HAIRLINE, w=1.5)
    node(s, hcx, hcy, hw, hh, "Coordinator", edge=RED, ew=1.6, bold=True, size=11, fill=SURFACE)
    for (x, y, t) in sats:
        node(s, x, y, 1.5, 0.46, t, fill=WHITE, size=9.5)


def fig_grounding(s):
    items = [("MIMIC-IV", HAIRLINE, False), ("Vector database", HAIRLINE, False),
             ("Patient-timeline context", RED, True), ("Clinical agents", HAIRLINE, False)]
    w = 3.0
    x = CXc - w / 2
    h = 0.5
    step = 0.735
    for i, (t, edge, emph) in enumerate(items):
        y = DY + i * step
        node(s, x, y, w, h, t, fill=(SURFACE if not emph else WHITE),
             edge=edge, ew=(1.6 if emph else 1.0), bold=emph, size=10.5)
        if i < 3:
            arrow(s, CXc, y + h, CXc, y + step, color=MUTE, w=1.3)


def fig_hitl(s):
    w = 2.6
    x = CXc - w / 2
    node(s, x, DY, w, 0.44, "Agent recommendation", size=9.5, fill=SURFACE)
    node(s, x, DY + 0.64, w, 0.44, "Verification & XAI", size=9.5, fill=SURFACE)
    dy = DY + 1.34
    diamond(s, CXc - 0.9, dy, 1.8, 0.78, "Clinician review", size=9)
    arrow(s, CXc, DY + 0.44, CXc, DY + 0.64)
    arrow(s, CXc, DY + 1.08, CXc, dy)
    outs = [(DX, "Approve", GREEN, GTINT), (CXc - 0.58, "Modify", ORANGE, OTINT),
            (DX + DW - 1.16, "Reject", RED, PTINT)]
    oy = dy + 1.06
    for (ox, t, col, tint) in outs:
        node(s, ox, oy, 1.16, 0.42, t, fill=tint, edge=col, tcolor=col, bold=True, size=9.5)
        arrow(s, CXc, dy + 0.78, ox + 0.58, oy, color=MUTE, w=1.2)


def fig_contributions(s):
    cards = [("Architectural", "Verifiable multi-agent framework"),
             ("Methodological", "Patient-timeline RAG + verifier"),
             ("Foundation", "Design + evaluation protocol")]
    w, h = 1.08, 1.48
    gap = (DW - 3 * w) / 2
    y = DY + 0.42
    for i, (t, d) in enumerate(cards):
        x = DX + i * (w + gap)
        node(s, x, y, w, h, "", fill=SURFACE, rad=0.12)
        badge(s, x + w / 2 - 0.17, y - 0.17, 0.34, color=GREEN)
        set_text(add_round(s, Inches(x), Inches(y + 0.12), Inches(w), Inches(h - 0.12), None),
                 [[(t, 10.5, INK, True)], [(" ", 3, INK, False)], [(d, 8.3, SLATE, False)]],
                 align=PP_ALIGN.CENTER, ls=1.05)


def fig_timeline(s):
    labels = ["Lit & taxonomy", "Framework", "Prototype & eval"]
    lab_w = 1.18
    track_x = DX + lab_w + 0.05
    track_w = DW - lab_w - 0.05
    ry0 = DY
    rh = 0.30
    rgap = 0.14
    bars = [(0.00, 0.42, ROSE), (0.30, 0.42, ROSE), (0.58, 0.42, RED)]
    for i, (lbl, (bx, bw, col)) in enumerate(zip(labels, bars)):
        y = ry0 + i * (rh + rgap)
        textbox(s, Inches(DX), Inches(y), Inches(lab_w), Inches(rh),
                [[(lbl, 8.5, INK, False, False)]], anchor=MSO_ANCHOR.MIDDLE)
        add_round(s, Inches(track_x + bx * track_w), Inches(y + 0.03),
                  Inches(bw * track_w), Inches(rh - 0.06), col, radius=0.4)
    axy = ry0 + 3 * (rh + rgap) - rgap + 0.04
    add_rect(s, Inches(track_x), Inches(axy), Inches(track_w), Pt(1), HAIRLINE)
    # scope box
    sy = axy + 0.20
    bw2 = (DW - 0.20) / 2
    for (bx, head, body, col, tint) in [
            (DX, "In scope", "Architecture · prototype · eval plan", GREEN, GTINT),
            (DX + bw2 + 0.20, "Out of scope", "Production · live EHR · trials", ORANGE, OTINT)]:
        node(s, bx, sy, bw2, 0.98, "", fill=WHITE, edge=HAIRLINE, rad=0.10)
        add_round(s, Inches(bx), Inches(sy), Inches(bw2), Inches(0.30), tint, radius=0.10)
        textbox(s, Inches(bx), Inches(sy + 0.02), Inches(bw2), Inches(0.28),
                [[(head, 9, col, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, Inches(bx + 0.08), Inches(sy + 0.36), Inches(bw2 - 0.16), Inches(0.58),
                [[(body, 8.5, SLATE, False, False)]], align=PP_ALIGN.CENTER, ls=1.1)


def fig_synthesis(s):
    ins = [(DX, DY + 0.15, "Agentic AI"), (DX, DY + 1.00, ["MIMIC-IV", "grounding"]),
           (DX, DY + 1.85, ["Trustworthy AI", "+ HITL"])]
    iw, ih = 1.6, 0.55
    ox, oy, ow, oh = DX + 2.05, DY + 0.72, 1.63, 1.02
    for (x, y, t) in ins:
        node(s, x, y, iw, ih, t, fill=SURFACE, size=9.5)
        arrow(s, x + iw, y + ih / 2, ox, oy + oh / 2)
    node(s, ox, oy, ow, oh, ["Intelligent", "clinical", "decision support"],
         edge=RED, ew=1.6, bold=True, size=10, fill=WHITE)


# ==========================================================================
# CONTENT  (section, title, [bullets], (fig-title, draw_fn))
# ==========================================================================
SLIDES = [
    ("BACKGROUND", "Clinical Data Overload",
     ["Unprecedented data generation: continuous streams of vitals, labs, notes, and imaging.",
      "Cognitive bottleneck: data production outpaces manual interpretation under time pressure.",
      "The role of AI: reduce the cognitive load of synthesis, not replace clinician judgment."],
     ("Clinical data bottleneck", fig_bottleneck)),
    ("BACKGROUND", "The Shift to Agentic AI",
     ["Passive LLMs: high fluency, but lack persistence, planning, and patient memory.",
      "Agentic AI: shifts models from passive responders to goal-directed systems.",
      "Core mechanisms: autonomous reasoning, tool use, memory management, environmental feedback."],
     ("Passive LLM vs. agentic loop", fig_passive_agentic)),
    ("PROBLEM STATEMENT", "Problem Statement",
     ["Volume & heterogeneity: reconciling diverse EHR data is slow and error-prone.",
      "LLM limitations: isolated tools lacking cross-step planning and durable patient models.",
      "CDSS deficits: high accuracy often lacks explainability, evidence retrieval, and long-term context."],
     ("Three failing pillars", fig_pillars)),
    ("LITERATURE REVIEW", "Agentic Frameworks",
     ["ReAct: interleaves reasoning and acting, but lacks long-term memory.",
      "AutoGen & CAMEL: enable multi-agent coordination and role-play, but struggle with communication loops.",
      "MetaGPT: enforces Standard Operating Procedures (SOPs), mirroring clinical protocols.",
      "Limitation: general-purpose designs lacking integrated clinical safety."],
     ("Framework capability matrix", fig_frameworks)),
    ("LITERATURE REVIEW", "RAG & Healthcare AI",
     ["Medical LLMs (e.g., Med-PaLM): strong QA performance but lack external retrieval and memory.",
      "Retrieval-Augmented Generation (RAG): anchors reasoning to verifiable evidence, reducing hallucinations.",
      "Current state (e.g., MedRAG): improves specificity but focuses on isolated diagnostic snapshots."],
     ("Retrieval-augmented generation", fig_rag)),
    ("RESEARCH GAP", "The Limits of Current Systems",
     ["Static evaluation: tested on isolated exams (e.g., MedQA) rather than real longitudinal trajectories.",
      "Generic retrieval: RAG grounds in general guidelines, missing the patient's specific evolving timeline.",
      "Omitted trustworthiness: verification and faithful audit trails are not first-class components."],
     ("The research gap", fig_venn)),
    ("RESEARCH OBJECTIVES", "Research Objectives",
     ["Primary: design an Agentic AI framework for patient monitoring and CDSS using MIMIC-IV.",
      "Secondary 1: integrate memory, RAG, and ReAct reasoning natively.",
      "Secondary 2: develop specialized clinical agents for monitoring, diagnosis, and treatment.",
      "Secondary 3: embed trustworthy-AI safeguards and human-in-the-loop validation architecturally."],
     ("Objective structure", fig_objectives)),
    ("RESEARCH QUESTIONS", "Research Questions",
     ["RQ1: What are the architectural limits of current LLM healthcare systems?",
      "RQ2: How can Agentic AI improve reasoning/memory while preserving physician control?",
      "RQ3: How can specialized agents collaborate while preserving explainability?",
      "RQ4: How does RAG improve reliability, and where does retrieval fail?",
      "RQ5: How can this framework be instantiated on the MIMIC-IV dataset?"],
     ("Five research questions", fig_rqs)),
    ("METHODOLOGY", "Proposed Methodology",
     ["Phase 1: critical literature review of LLMs, Agentic AI, and RAG.",
      "Phase 2: taxonomy formulation and research-gap analysis.",
      "Phase 3: design of the layered agentic framework architecture.",
      "Phase 4: bounded prototype specification & experimental evaluation plan."],
     ("Four-phase methodology", fig_methodology)),
    ("PROPOSED FRAMEWORK", "Architecture Overview",
     ["Perception layer: ingests EHR records, demographics, labs, notes.",
      "Memory layer: maintains short/long-term context and vector-based semantic stores.",
      "Reasoning & orchestration: sequences tasks via ReAct and routes among agents.",
      "Action & verification: generates recommendations under human review."],
     ("Layered architecture", fig_architecture)),
    ("PROPOSED FRAMEWORK", "Specialized Agent Roles",
     ["Coordinator agent: routes tasks and resolves conflicting recommendations.",
      "Monitoring agent: tracks physiological state over time.",
      "Diagnosis & treatment agents: weigh candidate conditions and propose interventions.",
      "Explanation & verification agents: summarize reasoning and check against guidelines."],
     ("Coordinator hub & spokes", fig_hub)),
    ("PROPOSED FRAMEWORK", "Grounding with RAG & MIMIC-IV",
     ["MIMIC-IV dataset: rich, retrospective critical-care records for longitudinal design.",
      "Patient-timeline RAG: shifts retrieval from generic guidelines to the patient's own evolving EHR.",
      "Context-aware decisions: fuses medical literature with patient-specific data for continuous monitoring."],
     ("Patient-timeline grounding", fig_grounding)),
    ("TRUSTWORTHY AI", "Trustworthy AI & Human-in-the-Loop",
     ["Explainable AI: surfaces exact evidence and reasoning for clinical review.",
      "Safety verification: algorithmic checks for contradictions and guideline compliance.",
      "Auditability: tamper-evident logging of retrieved data and agent interactions.",
      "Human-in-the-loop: clinicians retain final authority to approve, modify, or reject."],
     ("Human-in-the-loop workflow", fig_hitl)),
    ("CONTRIBUTIONS", "Expected Contributions",
     ["Architectural: a coherent, verifiable multi-agent framework for longitudinal CDSS.",
      "Methodological: patient-timeline retrieval paired with a dedicated verification agent.",
      "Foundation for future work: a reproducible design and concrete evaluation protocol."],
     ("Three contributions", fig_contributions)),
    ("SCOPE & TIMELINE", "Timeline & Scope Boundaries",
     ["In scope: conceptual architecture, bounded prototype, evaluation plan, MIMIC-IV integration.",
      "Out of scope: production hospital software, live EHR integration, prospective clinical trials.",
      "Timeline: Literature & Taxonomy → Framework Design → Prototype & Evaluation Planning."],
     ("Phases & scope", fig_timeline)),
]

CONCLUSION = ("CONCLUSION", "Conclusion & Discussion",
              ["Summary: transitioning from passive LLMs to verifiable, agent-based clinical workflows.",
               "Integration: fusing memory, RAG, and HITL enables safe, longitudinal monitoring.",
               "Final thought: transparent, evidence-based recommendations that assist — not replace — clinical judgment."],
              ("Synthesis", fig_synthesis))


# ==========================================================================
def logo_dims(width_in):
    from PIL import Image
    iw, ih = Image.open(LOGO).size
    w = Inches(width_in); h = Emu(int(w * ih / iw))
    return w, h


def add_footer(slide, page):
    add_rect(slide, Inches(0.7), Inches(6.98), Inches(11.93), Pt(1), HAIRLINE)
    textbox(slide, Inches(0.7), Inches(7.08), Inches(5.5), Inches(0.35),
            [[(SHORT_TITLE, 10.5, SLATE, False, False)]])
    textbox(slide, Inches(4.0), Inches(7.08), Inches(5.33), Inches(0.35),
            [[(UNIV, 10.5, SLATE, False, False)]], align=PP_ALIGN.CENTER)
    textbox(slide, Inches(11.6), Inches(7.08), Inches(1.03), Inches(0.35),
            [[(str(page), 10.5, SLATE, False, False)]], align=PP_ALIGN.RIGHT)


def content_slide(prs, blank, section, title, bullets, fig, page, fig_no):
    fig_title, draw_fn = fig
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    if os.path.exists(LOGO):
        lw, lh = logo_dims(1.35)
        s.shapes.add_picture(LOGO, int(SLIDE_W - lw - Inches(0.7)), Inches(0.5), width=lw, height=lh)
    textbox(s, Inches(0.9), Inches(0.55), Inches(6.5), Inches(0.3),
            [[(section, 12.5, RED, False, False)]])
    add_rect(s, Inches(0.7), Inches(0.92), Pt(4), Inches(0.7), RED)
    textbox(s, Inches(0.9), Inches(0.9), Inches(7.2), Inches(0.85),
            [[(title, 33, INK, False, False)]], anchor=MSO_ANCHOR.MIDDLE)
    # body
    body = [[("—  ", 19, RED, False, False), (b, 19, INK, False, False)] for b in bullets]
    textbox(s, Inches(0.9), Inches(2.05), Inches(7.1), Inches(4.6),
            body, space_after=14, ls=1.12)
    # figure card + header
    add_round(s, Inches(CARD_X), Inches(CARD_Y), Inches(CARD_W), Inches(CARD_H),
              SURFACE, radius=0.05, shadow=True)
    textbox(s, Inches(CARD_X + 0.30), Inches(CARD_Y + 0.26), Inches(CARD_W - 0.6), Inches(0.26),
            [[("FIGURE %d" % fig_no, 11, MUTE, False, False)]])
    textbox(s, Inches(CARD_X + 0.30), Inches(CARD_Y + 0.52), Inches(CARD_W - 0.6), Inches(0.34),
            [[(fig_title, 13.5, INK, False, False)]])
    add_rect(s, Inches(CARD_X + 0.30), Inches(CARD_Y + 0.90), Inches(0.55), Pt(2.5), ROSE)
    draw_fn(s)
    add_footer(s, page)


def build(path):
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    blank = prs.slide_layouts[6]

    # Title slide (dark canvas)
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, DARK)
    if os.path.exists(LOGO):
        lw, lh = logo_dims(3.4)
        s.shapes.add_picture(LOGO, int((SLIDE_W - lw) / 2), Inches(0.85), width=lw, height=lh)
    add_rect(s, int((SLIDE_W - Inches(2.6)) / 2), Inches(3.35), Inches(2.6), Pt(3), RED)
    textbox(s, Inches(1.2), Inches(3.65), SLIDE_W - Inches(2.4), Inches(1.7),
            [[("An Agentic AI Framework for Intelligent Patient", 31, WHITE, False, False)],
             [("Monitoring and Clinical Decision Support", 31, WHITE, False, False)]],
            align=PP_ALIGN.CENTER, ls=1.1)
    textbox(s, Inches(1.2), Inches(5.5), SLIDE_W - Inches(2.4), Inches(1.6),
            [[("Adeel Gill", 19, WHITE, False, False)],
             [("MS Thesis Proposal Defense  ·  Master of Science in Artificial Intelligence", 14, DARKSUB, False, False)],
             [("Supervisor: Dr. Fawad Nasim", 14, DARKSUB, False, False)],
             [("Faculty of Computer Science and Information Technology  ·  The Superior University, Lahore", 13, DARKSUB, False, False)]],
            align=PP_ALIGN.CENTER, space_after=4, ls=1.05)

    # Content slides
    page = 2; fno = 1
    for (section, title, bullets, fig) in SLIDES:
        content_slide(prs, blank, section, title, bullets, fig, page, fno)
        page += 1; fno += 1
    content_slide(prs, blank, *CONCLUSION[:3], CONCLUSION[3], page, fno)
    page += 1

    # Thank You (dark bookend)
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, DARK)
    add_rect(s, int((SLIDE_W - Inches(2.6)) / 2), Inches(2.75), Inches(2.6), Pt(3), RED)
    textbox(s, Inches(1.2), Inches(3.05), SLIDE_W - Inches(2.4), Inches(1.2),
            [[("Thank You", 44, WHITE, False, False)]], align=PP_ALIGN.CENTER)
    textbox(s, Inches(1.2), Inches(4.25), SLIDE_W - Inches(2.4), Inches(0.6),
            [[("Questions & Discussion", 20, DARKSUB, False, False)]], align=PP_ALIGN.CENTER)
    textbox(s, Inches(1.2), Inches(5.15), SLIDE_W - Inches(2.4), Inches(0.6),
            [[("Adeel Gill  ·  adeel@innovahealth.com  ·  The Superior University, Lahore",
               13, DARKSUB, False, False)]], align=PP_ALIGN.CENTER)

    prs.save(path)
    print("PPTX written:", path, "(%d slides)" % len(prs.slides._sldIdLst))


if __name__ == "__main__":
    build(os.path.join(HERE, "Defense_Presentation_v1.pptx"))
