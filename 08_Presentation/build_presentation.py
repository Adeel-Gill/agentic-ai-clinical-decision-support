#!/usr/bin/env python3
"""
Build the defense deliverables:
  * Defense_Presentation.pptx  -- 17-slide PowerPoint deck (from Defense_Presentation.md)
  * Defense_Speaker_Notes.pdf  -- formatted speaker-notes booklet (from Defense_Speaker_Notes.md)

Run:  python build_presentation.py
Requires: python-pptx, reportlab  (logo PNG at ../07_Thesis/Images/superior_logo.png)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
LOGO = os.path.join(HERE, "..", "07_Thesis", "Images", "superior_logo.png")

# ---- Brand palette -------------------------------------------------------
PURPLE      = RGBColor(0x70, 0x1A, 0x73)   # Superior University brand
DARK_BG     = RGBColor(0x2A, 0x14, 0x2E)   # deep purple, title/section slides
LIGHT_TXT   = RGBColor(0xFF, 0xFF, 0xFF)
BODY_TXT    = RGBColor(0x2B, 0x2B, 0x2B)
ACCENT      = RGBColor(0x9C, 0x27, 0xB0)   # lighter accent
GREY_TXT    = RGBColor(0x6B, 0x6B, 0x6B)
FONT        = "Calibri"

# ==========================================================================
# SLIDE CONTENT  (title, [bullets], visual-note)
# ==========================================================================
SLIDES = [
    # 2
    ("Background: Clinical Data Overload",
     ["Unprecedented data generation: continuous streams of vitals, labs, notes, and imaging.",
      "Cognitive bottleneck: data production outpaces manual interpretation under time pressure.",
      "The role of AI: reduce the cognitive load of synthesis, not replace clinician judgment."],
     "Visual: ICU data streams funneling through a narrow bottleneck to a single clinician."),
    # 3
    ("Background: The Shift to Agentic AI",
     ["Passive LLMs: high fluency, but lack persistence, planning, and patient memory.",
      "Agentic AI: shifts models from passive responders to goal-directed systems.",
      "Core mechanisms: autonomous reasoning, tool use, memory management, environmental feedback."],
     "Visual: Passive LLM (Prompt -> Black Box -> Response) vs. Agentic Loop (Perceive-Memory-Reason-Act)."),
    # 4
    ("Problem Statement",
     ["Volume & heterogeneity: reconciling diverse EHR data is slow and error-prone.",
      "LLM limitations: isolated tools lacking cross-step planning and durable patient models.",
      "CDSS deficits: high accuracy often lacks explainability, evidence retrieval, and long-term context."],
     "Visual: Three crumbling pillars - Lack of Context, Opaque Reasoning, Fragmented Systems."),
    # 5
    ("Literature Review: Agentic Frameworks",
     ["ReAct: interleaves reasoning and acting, but lacks long-term memory.",
      "AutoGen & CAMEL: enable multi-agent coordination and role-play, but struggle with communication loops.",
      "MetaGPT: enforces Standard Operating Procedures (SOPs), mirroring clinical protocols.",
      "Limitation: general-purpose designs lacking integrated clinical safety."],
     "Visual: Comparison matrix - frameworks vs. clinical needs (Memory, SOPs, Oversight)."),
    # 6
    ("Literature Review: RAG & Healthcare AI",
     ["Medical LLMs (e.g., Med-PaLM): strong QA performance but lack external retrieval and memory.",
      "Retrieval-Augmented Generation (RAG): anchors reasoning to verifiable evidence, reducing hallucinations.",
      "Current state (e.g., MedRAG): improves specificity but focuses on isolated diagnostic snapshots."],
     "Visual: Generator LLM fed by a Retriever over External Medical Knowledge + EHR Databases."),
    # 7
    ("Research Gap: The Limits of Current Systems",
     ["Static evaluation: tested on isolated exams (e.g., MedQA) rather than real longitudinal trajectories.",
      "Generic retrieval: RAG grounds in general guidelines, missing the patient's specific evolving timeline.",
      "Omitted trustworthiness: verification and faithful audit trails are not first-class components."],
     "Visual: Venn of Longitudinal ICU Evaluation, Patient-Timeline RAG, Faithful Auditability -> The Missing Framework."),
    # 8
    ("Research Objectives",
     ["Primary: design an Agentic AI framework for patient monitoring and CDSS using MIMIC-IV.",
      "Secondary 1: integrate memory, RAG, and ReAct reasoning natively.",
      "Secondary 2: develop specialized clinical agents for monitoring, diagnosis, and treatment.",
      "Secondary 3: embed trustworthy-AI safeguards and human-in-the-loop validation architecturally."],
     "Visual: Target with one central objective and three supporting nodes."),
    # 9
    ("Research Questions",
     ["RQ1: What are the architectural limits of current LLM healthcare systems?",
      "RQ2: How can Agentic AI improve reasoning/memory while preserving physician control?",
      "RQ3: How can specialized agents collaborate while preserving explainability?",
      "RQ4: How does RAG improve reliability, and where does retrieval fail?",
      "RQ5: How can this framework be instantiated on the MIMIC-IV dataset?"],
     "Visual: Five sequential chevrons leading to a blueprint icon."),
    # 10
    ("Proposed Methodology",
     ["Phase 1: critical literature review of LLMs, Agentic AI, and RAG.",
      "Phase 2: taxonomy formulation and research-gap analysis.",
      "Phase 3: design of the layered agentic framework architecture.",
      "Phase 4: bounded prototype specification & experimental evaluation plan."],
     "Visual: Four-step process graphic, Phase 1 -> Phase 4."),
    # 11
    ("Proposed Framework: Architecture Overview",
     ["Perception layer: ingests EHR records, demographics, labs, notes.",
      "Memory layer: maintains short/long-term context and vector-based semantic stores.",
      "Reasoning & orchestration: sequences tasks via ReAct and routes among agents.",
      "Action & verification: generates recommendations under human review."],
     "Visual: Stacked layers Perception -> Memory -> Reasoning -> Action, with a vertical Trust pillar."),
    # 12
    ("Proposed Framework: Specialized Agent Roles",
     ["Coordinator agent: routes tasks and resolves conflicting recommendations.",
      "Monitoring agent: tracks physiological state over time.",
      "Diagnosis & treatment agents: weigh candidate conditions and propose interventions.",
      "Explanation & verification agents: summarize reasoning and check against guidelines."],
     "Visual: Hub-and-spoke - Coordinator routing to specialized agents."),
    # 13
    ("Grounding the Framework: RAG & MIMIC-IV",
     ["MIMIC-IV dataset: rich, retrospective critical-care records for longitudinal design.",
      "Patient-timeline RAG: shifts retrieval from generic guidelines to the patient's own evolving EHR.",
      "Context-aware decisions: fuses medical literature with patient-specific data for continuous monitoring."],
     "Visual: MIMIC-IV -> Vector Database -> patient-specific timeline context to agents."),
    # 14
    ("Trustworthy AI & Human-in-the-Loop",
     ["Explainable AI: surfaces exact evidence and reasoning for clinical review.",
      "Safety verification: algorithmic checks for contradictions and guideline compliance.",
      "Auditability: tamper-evident logging of retrieved data and agent interactions.",
      "Human-in-the-loop: clinicians retain final authority to approve, modify, or reject."],
     "Visual: Workflow ending in an Approve / Modify / Reject decision node."),
    # 15
    ("Expected Contributions",
     ["Architectural: a coherent, verifiable multi-agent framework for longitudinal CDSS.",
      "Methodological: patient-timeline retrieval paired with a dedicated verification agent.",
      "Foundation for future work: a reproducible design and concrete evaluation protocol."],
     "Visual: Three pillars with checkmark icons."),
    # 16
    ("Timeline & Scope Boundaries",
     ["In scope: conceptual architecture, bounded prototype, evaluation plan, MIMIC-IV integration.",
      "Out of scope: production hospital software, live EHR integration, prospective clinical trials.",
      "Timeline: Literature & Taxonomy -> Framework Design -> Prototype & Evaluation Planning."],
     "Visual: Gantt of thesis phases plus an In-Scope / Out-of-Scope box."),
    # 17
    ("Conclusion & Q&A",
     ["Summary: transitioning from passive LLMs to verifiable, agent-based clinical workflows.",
      "Integration: fusing memory, RAG, and HITL enables safe, longitudinal monitoring.",
      "Final thought: transparent, evidence-based recommendations that assist - not replace - clinical judgment."],
     "Visual: Agentic AI + MIMIC-IV + Trustworthy AI  ->  Intelligent Clinical Decision Support."),
]

# ==========================================================================
def build_pptx(path):
    prs = Presentation()
    prs.slide_width  = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def add_rect(slide, x, y, w, h, color):
        from pptx.enum.shapes import MSO_SHAPE
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid(); shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    # ---------- Title slide ----------
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, DARK_BG)
    # logo (centered, top)
    if os.path.exists(LOGO):
        from PIL import Image
        iw, ih = Image.open(LOGO).size
        disp_w = Inches(3.6)
        disp_h = Emu(int(disp_w * ih / iw))
        s.shapes.add_picture(LOGO, int((SW - disp_w) / 2), Inches(0.7),
                             width=disp_w, height=disp_h)
    # accent line
    add_rect(s, int((SW - Inches(3.0)) / 2), Inches(3.35), Inches(3.0), Pt(3), ACCENT)
    # title
    tb = s.shapes.add_textbox(Inches(1.0), Inches(3.6), SW - Inches(2.0), Inches(1.7))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = ("An Agentic AI Framework for Intelligent Patient "
                               "Monitoring and Clinical Decision Support")
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = LIGHT_TXT; r.font.name = FONT
    # meta
    meta = [("Adeel Gill", 20, True),
            ("MS Thesis Proposal Defense  •  Master of Science in Artificial Intelligence", 15, False),
            ("Supervisor: Dr. Fawad Nasim", 15, False),
            ("Faculty of Computer Science and Information Technology", 13, False),
            ("The Superior University, Lahore", 13, False)]
    mb = s.shapes.add_textbox(Inches(1.0), Inches(5.4), SW - Inches(2.0), Inches(1.8))
    mf = mb.text_frame; mf.word_wrap = True
    for i, (txt, sz, bold) in enumerate(meta):
        para = mf.paragraphs[0] if i == 0 else mf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run(); run.text = txt
        run.font.size = Pt(sz); run.font.bold = bold
        run.font.color.rgb = LIGHT_TXT if bold else RGBColor(0xD9, 0xC7, 0xDD)
        run.font.name = FONT
        para.space_after = Pt(2)

    # ---------- Content slides ----------
    for idx, (title, bullets, visual) in enumerate(SLIDES, start=2):
        s = prs.slides.add_slide(blank)
        add_rect(s, 0, 0, SW, SH, RGBColor(0xFF, 0xFF, 0xFF))
        # header band
        add_rect(s, 0, 0, SW, Inches(1.15), PURPLE)
        # small logo top-right
        if os.path.exists(LOGO):
            from PIL import Image
            iw, ih = Image.open(LOGO).size
            lw = Inches(1.5); lh = Emu(int(lw * ih / iw))
            # white plate behind logo for contrast
            add_rect(s, SW - lw - Inches(0.35), Inches(0.28), lw + Inches(0.2), lh + Inches(0.12),
                     RGBColor(0xFF, 0xFF, 0xFF))
            s.shapes.add_picture(LOGO, int(SW - lw - Inches(0.25)), Inches(0.34), width=lw, height=lh)
        # title text
        tb = s.shapes.add_textbox(Inches(0.55), Inches(0.12), SW - Inches(2.6), Inches(0.95))
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = title
        r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = LIGHT_TXT; r.font.name = FONT
        # bullets
        bb = s.shapes.add_textbox(Inches(0.7), Inches(1.55), SW - Inches(1.4), Inches(4.6))
        bf = bb.text_frame; bf.word_wrap = True
        for i, b in enumerate(bullets):
            para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
            run = para.add_run(); run.text = "▸  " + b
            run.font.size = Pt(19); run.font.color.rgb = BODY_TXT; run.font.name = FONT
            para.space_after = Pt(14); para.line_spacing = 1.1
        # visual-recommendation footnote
        vb = s.shapes.add_textbox(Inches(0.7), Inches(6.35), SW - Inches(1.4), Inches(0.8))
        vf = vb.text_frame; vf.word_wrap = True
        vp = vf.paragraphs[0]
        vr = vp.add_run(); vr.text = visual
        vr.font.size = Pt(11); vr.font.italic = True; vr.font.color.rgb = GREY_TXT; vr.font.name = FONT
        # slide number
        nb = s.shapes.add_textbox(SW - Inches(1.0), SH - Inches(0.45), Inches(0.7), Inches(0.35))
        nf = nb.text_frame
        nr = nf.paragraphs[0].add_run(); nr.text = str(idx)
        nr.font.size = Pt(11); nr.font.color.rgb = GREY_TXT; nr.font.name = FONT
        nf.paragraphs[0].alignment = PP_ALIGN.RIGHT

    prs.save(path)
    print("PPTX written:", path, "(%d slides)" % len(prs.slides._sldIdLst))


# ==========================================================================
# Speaker-notes PDF
# ==========================================================================
def build_notes_pdf(md_path, pdf_path):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm
    from reportlab.lib.enums import TA_LEFT, TA_CENTER
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.colors import HexColor
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Image as RLImage,
                                    HRFlowable)
    import re

    purple = HexColor("#701A73")
    grey   = HexColor("#6B6B6B")

    styles = getSampleStyleSheet()
    h_title = ParagraphStyle("hTitle", parent=styles["Title"], textColor=purple,
                             fontSize=20, leading=24, alignment=TA_CENTER, spaceAfter=4)
    h_sub   = ParagraphStyle("hSub", parent=styles["Normal"], textColor=grey,
                             fontSize=10, leading=14, alignment=TA_CENTER, spaceAfter=2)
    h_slide = ParagraphStyle("hSlide", parent=styles["Heading2"], textColor=purple,
                             fontSize=13, leading=16, spaceBefore=14, spaceAfter=4)
    body    = ParagraphStyle("body", parent=styles["Normal"], fontSize=11, leading=16,
                             alignment=TA_LEFT, spaceAfter=6)
    pacing  = ParagraphStyle("pacing", parent=styles["Normal"], fontSize=10.5, leading=15,
                             textColor=HexColor("#333333"), backColor=HexColor("#F3E9F5"),
                             borderPadding=6, spaceAfter=10)

    def md_inline(t):
        t = t.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        t = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", t)
        t = re.sub(r"(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)", r"<i>\1</i>", t)
        t = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", t)   # strip md links -> text
        t = t.replace("~", "&#126;")
        return t

    story = []
    if os.path.exists(LOGO):
        from PIL import Image as PILImage
        iw, ih = PILImage.open(LOGO).size
        w = 55 * mm; h = w * ih / iw
        img = RLImage(LOGO, width=w, height=h); img.hAlign = "CENTER"
        story += [img, Spacer(1, 6)]
    story.append(Paragraph("Defense &mdash; Speaker Notes", h_title))
    story.append(Paragraph("An Agentic AI Framework for Intelligent Patient Monitoring "
                           "and Clinical Decision Support", h_sub))
    story.append(Paragraph("Adeel Gill &nbsp;&bull;&nbsp; Supervisor: Dr. Fawad Nasim &nbsp;&bull;&nbsp; "
                           "The Superior University, Lahore", h_sub))
    story.append(Spacer(1, 6))
    story.append(HRFlowable(width="100%", thickness=1, color=purple))
    story.append(Spacer(1, 4))

    with open(md_path, encoding="utf-8") as f:
        lines = f.read().splitlines()

    # skip the markdown H1 title block; render from first pacing/### onward
    buf = []
    def flush():
        if buf:
            story.append(Paragraph(md_inline(" ".join(buf).strip()), body))
            buf.clear()

    in_body = False
    for ln in lines:
        s = ln.strip()
        if s.startswith("# "):            # H1 -> already have our own title
            continue
        if s.startswith("### "):
            flush(); in_body = True
            story.append(Paragraph(md_inline(s[4:]), h_slide))
            continue
        if s.startswith("**Pacing"):
            flush(); in_body = True
            story.append(Paragraph(md_inline(s), pacing))
            continue
        if s.startswith("---"):
            continue
        if not in_body:
            # intro paragraph(s) before first heading
            if s == "":
                flush()
            else:
                buf.append(s)
            continue
        if s == "":
            flush()
        else:
            buf.append(s)
    flush()

    doc = SimpleDocTemplate(pdf_path, pagesize=A4,
                            leftMargin=20*mm, rightMargin=20*mm,
                            topMargin=16*mm, bottomMargin=16*mm,
                            title="Defense Speaker Notes", author="Adeel Gill")
    doc.build(story)
    print("PDF written:", pdf_path)


if __name__ == "__main__":
    build_pptx(os.path.join(HERE, "Defense_Presentation.pptx"))
    build_notes_pdf(os.path.join(HERE, "Defense_Speaker_Notes.md"),
                    os.path.join(HERE, "Defense_Speaker_Notes.pdf"))
